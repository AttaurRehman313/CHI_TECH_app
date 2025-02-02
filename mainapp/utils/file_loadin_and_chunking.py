from langchain_text_splitters import RecursiveCharacterTextSplitter
from mainapp.model.image_to_text import image_text
from langchain.schema import Document
from pptx import Presentation
import os,shutil
from time import time,sleep


text_runs = []   # for text
image_text_runs = []  # for images text
image_runs=[] # for images

image_count=len(image_runs)

def extract_text_and_image_from_slide(pptx_file):

    # Load the presentation
    prs = Presentation(pptx_file)

    # Create the output folder for images, if it doesn't exist
    if not os.path.exists("extracted_images"):
        os.makedirs("extracted_images")
    else:
        shutil.rmtree("extracted_images")
        os.makedirs("extracted_images")

    # Loop through slides
    for slide_idx, slide in enumerate(prs.slides, start=1):
        for shape_idx, shape in enumerate(slide.shapes, start=1):# Loop through shapes on the slide
            if shape.has_text_frame:# Extract text from text boxes
                for paragraph in shape.text_frame.paragraphs:
                    for run in paragraph.runs:
                        text_runs.append(run.text)

            # Extract images from shapes
            if shape.shape_type == 13:  # 13 corresponds to a Picture shape
                image = shape.image
                image_bytes = image.blob  # Image data in bytes
                image_format = image.ext  # Image file extension (e.g., 'png', 'jpeg')


                # Generate a unique file name for the image
                image_filename = f"slide{slide_idx}_shape{shape_idx}.{image_format}"
                image_path = os.path.join("extracted_images", image_filename)

                # Save the image to the output folder
                with open(image_path, "wb") as img_file:
                    img_file.write(image_bytes)
                
                # Store image details in list
                image_runs.append(image_path)


# extract_text_and_image_from_slide(pptx_file)



# # Print extracted text and image paths
# print("Extracted Text:")
# for text in text_runs:
#     print(text)

# print("\nExtracted Images:************************************\n")
# for image_path in image_runs:
#     print(image_path)





start_time = time()
def transcribe_image(image_runs):

    for index, image_path in enumerate(image_runs):
        if index==0:
            try:
                text_of_image = image_text(image_path )
                image_text_runs.append(text_of_image)
                print(f"Index: {image_runs.index(image_path)}")
            except Exception as e:
                print(f"Error generating caption for {image_path}: {e}")

        elif index%10!=0:
            try:
                text_of_image = image_text(image_path )
                image_text_runs.append(text_of_image)
                print(f"Index: {image_runs.index(image_path)}")

            except Exception as e:
                print(f"Error generating caption for {image_path}: {e}")
        else:
            sleep(21)
            try:
                text_of_image = image_text(image_path )
                image_text_runs.append(text_of_image)
                print(f"Index: {image_runs.index(image_path)}")

            except Exception as e:
                print(f"Error generating caption for {image_path}: {e}")


    combined_text = text_runs + image_text_runs    
    data = "\n\n".join(combined_text)
    print("\n\n\n\ncombined data \n",data)
    return data


end_time = time()
elapsed_time = end_time - start_time
elapsed_minutes = elapsed_time / 60
print(f"************************************* Elapsed time: {elapsed_time:.2f} seconds*********************************\n")




# def combine_lists_data():
#     combined_text = text_runs + image_text_runs    
#     data = "\n\n".join(combined_text)
#     print("\n\n\n\ncombined data \n",data)
#     return data

# data=combine_lists_data()

# data=transcribe_image(image_runs)


def create_chunks(data):
    """
    This function loads the provided file data and splits it into chunks 
    using RecursiveCharacterTextSplitter with a specified chunk size 
    and overlap.
    import os
    """
    text=data
    text_splitter = RecursiveCharacterTextSplitter(chunk_size=100, chunk_overlap=10,separators=["\n\n"])
    text_chunks = text_splitter.split_text(text)
    documents = [
    Document(
        metadata={'source': 'output.pdf', 'page': i, 'page_label': str(i + 1)},
        page_content=f"Slide Content:\n{chunk}"
    )
    for i, chunk in enumerate(text_chunks)
    ]
    return documents

# documents=create_chunks(data)


# documents = [
#     Document(
#         metadata={'source': 'output.pdf', 'page': i, 'page_label': str(i + 1)},
#         page_content=f"Slide Content:\n{chunk}"
#     )
#     for i, chunk in enumerate(chunks)
# ]
# print("\n\n\n\n cunks\n",documents)








